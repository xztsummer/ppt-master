#!/usr/bin/env python3
"""
PPT Master - Template Preview PPTX Exporter

Export complete Slide SVG prototypes as a structured review deck.

Usage:
    python3 scripts/template_preview_pptx.py <template_workspace> [-o output.pptx]

Examples:
    python3 scripts/template_preview_pptx.py projects/my_template
    python3 scripts/template_preview_pptx.py templates/decks/my_template -o review.pptx

Dependencies:
    python-pptx
"""

from __future__ import annotations

import argparse
import contextlib
import re
import shutil
import sys
import tempfile
from collections.abc import Iterator
from pathlib import Path
from xml.etree import ElementTree as ET

from attribution_guard import require_skill_integrity
from console_encoding import configure_utf8_stdio


configure_utf8_stdio()

from pptx import Presentation  # noqa: E402

from svg_to_pptx.drawingml.theme_fonts import (  # noqa: E402
    infer_master_text_style_spec,
)
from svg_to_pptx.pptx_package.builder import (  # noqa: E402
    create_pptx_with_native_svg,
)
from svg_to_pptx.pptx_package.template_structure import (  # noqa: E402
    load_template_source_themes,
)


_FRONTMATTER_ID_RE = re.compile(
    r"^(?:template_id|deck_id|layout_id)\s*:\s*(.+?)\s*$",
    re.MULTILINE,
)
_REPLICATION_MODE_RE = re.compile(
    r"^replication_mode\s*:\s*(standard|fidelity|mirror)\s*$",
    re.MULTILINE,
)
_CANVAS_VIEWBOX_RE = re.compile(
    r"^canvas_viewbox\s*:\s*[\"']?([^\"'\r\n]+?)[\"']?\s*$",
    re.MULTILINE,
)
_FILENAME_UNSAFE_RE = re.compile(r"[\\/:*?\"<>|\x00-\x1f]+")
_PLACEHOLDER_MARKER_RE = re.compile(r"\{\{([A-Z][A-Z0-9_]*)\}\}")


def _review_marker_text(match: re.Match[str]) -> str:
    """Return concise preview-only text for one canonical marker."""
    token = match.group(1)
    if token in {"PAGE_NUM", "SLIDE_NUM"}:
        return "1"
    if token.endswith("_NUM"):
        return "01"
    if token == "DATE":
        return "YYYY-MM-DD"
    return token.replace("_", " ").title()


def _write_review_svg(source: Path, target: Path) -> bool:
    """Copy one SVG, shortening only visible placeholder-carrier prompts."""
    tree = ET.parse(source)
    changed = False
    for slot in tree.getroot().iter():
        if not (slot.get("data-pptx-placeholder") or "").strip():
            continue
        for carrier in slot.iter():
            if (
                carrier.get("data-pptx-carrier") or ""
            ).strip().lower() != "true":
                continue
            for element in carrier.iter():
                if element.text:
                    updated = _PLACEHOLDER_MARKER_RE.sub(
                        _review_marker_text,
                        element.text,
                    )
                    if updated != element.text:
                        element.text = updated
                        changed = True
    # The staging directory becomes the converter's project root, so
    # workspace-relative bitmap hrefs must be rebased into the staging copy
    # or the resolver's root boundary filters them out.
    for image in tree.getroot().iter():
        if image.tag != "{http://www.w3.org/2000/svg}image":
            continue
        href = image.get("href") or ""
        if href.startswith("../images/"):
            image.set("href", "images/" + href[len("../images/"):])
            changed = True
    if changed:
        tree.write(target, encoding="utf-8", xml_declaration=True)
    else:
        shutil.copy2(source, target)
    return changed


@contextlib.contextmanager
def _review_svg_sources(
    workspace: Path,
    svg_files: list[Path],
    *,
    shorten_placeholder_markers: bool,
) -> Iterator[list[Path]]:
    """Yield ephemeral review SVGs without modifying canonical template files."""
    if not shorten_placeholder_markers:
        yield svg_files
        return

    with tempfile.TemporaryDirectory(
        prefix=".template-preview-",
        dir=workspace,
    ) as temporary:
        review_dir = Path(temporary)
        images_dir = workspace / "images"
        if images_dir.is_dir():
            shutil.copytree(images_dir, review_dir / "images", dirs_exist_ok=True)
        review_files: list[Path] = []
        shortened = 0
        for source in svg_files:
            target = review_dir / source.name
            shortened += int(_write_review_svg(source, target))
            review_files.append(target)
        print(
            "  Review prompt text: preview-only samples in "
            f"{shortened} SVG(s); canonical {{{{...}}}} markers unchanged"
        )
        yield review_files


_TEMPLATE_SPEC_NAME_RE = re.compile(
    r"design_spec\.(?P<kind>brand|style|layout|deck)\.[^/\\]+\.md"
)


def _roster_spec(directory: Path) -> Path | None:
    """Return the effective spec that owns this directory's SVG roster.

    A library workspace keeps the exact ``design_spec.md``. A project workspace
    shares one ``templates/`` across kinds. Layout owns structure when both
    Layout and Deck are present; otherwise Deck owns it.
    """
    if not directory.is_dir():
        return None
    exact = directory / "design_spec.md"
    qualified = []
    for item in sorted(directory.glob("design_spec.*.md")):
        match = _TEMPLATE_SPEC_NAME_RE.fullmatch(item.name)
        if match is not None:
            qualified.append((item, match.group("kind")))
    if exact.is_file() and qualified:
        raise ValueError(
            "design_spec.md and design_spec.<kind>.<id>.md cannot share "
            f"{directory}; rename the bare spec to its kind-qualified name"
        )
    kinds = [kind for _item, kind in qualified]
    duplicate_kinds = sorted({
        kind for kind in kinds if kinds.count(kind) > 1
    })
    if duplicate_kinds:
        raise ValueError(
            f"{directory} declares the same kind more than once: "
            + ", ".join(duplicate_kinds)
        )
    try:
        from register_template import (
            SpecParseError,
            validate_qualified_spec_identity,
        )
        for item, _kind in qualified:
            validate_qualified_spec_identity(item)
    except ImportError as exc:
        raise ValueError(
            f"Qualified Design Spec validator could not be imported: {exc}"
        ) from exc
    except (OSError, SpecParseError) as exc:
        raise ValueError(str(exc)) from exc
    if exact.is_file():
        return exact
    for preferred_kind in ("layout", "deck"):
        for item, kind in qualified:
            if kind == preferred_kind:
                return item
    return None


def _resolve_workspace(path: Path) -> tuple[Path, Path]:
    """Resolve one workspace root and its canonical template-source directory."""
    candidate = path.expanduser().resolve()
    if _roster_spec(candidate / "templates") is not None:
        return candidate, candidate / "templates"
    raise ValueError(
        "template workspace root must contain templates/design_spec.md or "
        "templates/design_spec.<layout|deck>.<id>.md"
    )


def _template_id(spec_path: Path) -> str:
    """Read the required portable template id."""
    text = spec_path.read_text(encoding="utf-8")
    match = _FRONTMATTER_ID_RE.search(text)
    if match is None:
        raise ValueError(
            f"{spec_path.name} frontmatter must declare layout_id or deck_id"
        )
    raw = match.group(1).strip().strip("'\"")
    safe = _FILENAME_UNSAFE_RE.sub("_", raw).strip(" ._")
    return safe or "template"


def _replication_mode(spec_path: Path) -> str:
    """Read the required template replication mode."""
    text = spec_path.read_text(encoding="utf-8")
    match = _REPLICATION_MODE_RE.search(text)
    if match is None:
        raise ValueError(
            f"{spec_path.name} frontmatter must declare replication_mode"
        )
    return match.group(1)


def _canvas_viewbox(spec_path: Path) -> str | None:
    """Read the template's locked root canvas when declared."""
    text = spec_path.read_text(encoding="utf-8")
    if not text.startswith("---\n"):
        return None
    end = text.find("\n---\n", 4)
    if end == -1:
        return None
    match = _CANVAS_VIEWBOX_RE.search(text[4:end])
    return match.group(1).strip() if match else None


def _verify_output(
    output_path: Path,
    *,
    require_full_placeholder_frames: bool,
) -> tuple[int, int, int, int]:
    """Reopen the review deck and verify counts plus authored placeholder frames."""
    presentation = Presentation(str(output_path))
    master_count = len(presentation.slide_masters)
    layout_count = sum(len(master.slide_layouts) for master in presentation.slide_masters)
    placeholder_count = 0
    if require_full_placeholder_frames:
        for slide_number, slide in enumerate(presentation.slides, 1):
            layout_placeholders = {
                shape.placeholder_format.idx: shape
                for shape in slide.slide_layout.placeholders
            }
            slide_placeholders = {
                shape.placeholder_format.idx: shape
                for shape in slide.placeholders
            }
            if set(slide_placeholders) != set(layout_placeholders):
                raise ValueError(
                    f"review slide {slide_number} placeholder indexes do not match "
                    f"its Layout: {sorted(slide_placeholders)} != "
                    f"{sorted(layout_placeholders)}"
                )
            for placeholder_idx, slide_shape in slide_placeholders.items():
                layout_shape = layout_placeholders[placeholder_idx]
                if (
                    slide_shape.placeholder_format.type
                    != layout_shape.placeholder_format.type
                ):
                    raise ValueError(
                        f"review slide {slide_number} placeholder {placeholder_idx} "
                        "type does not match its Layout"
                    )
                slide_frame = (
                    slide_shape.left,
                    slide_shape.top,
                    slide_shape.width,
                    slide_shape.height,
                )
                layout_frame = (
                    layout_shape.left,
                    layout_shape.top,
                    layout_shape.width,
                    layout_shape.height,
                )
                if slide_frame != layout_frame:
                    raise ValueError(
                        f"review slide {slide_number} placeholder {placeholder_idx} "
                        f"uses a tight/local frame {slide_frame}; expected full "
                        f"Layout frame {layout_frame}"
                    )
                placeholder_count += 1
    return len(presentation.slides), master_count, layout_count, placeholder_count


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Export a complete template workspace as a structured PPTX review deck."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "template_workspace",
        help="Workspace root containing templates/design_spec.md.",
    )
    parser.add_argument(
        "-o",
        "--output",
        help=(
            "Output PPTX path. Default: "
            "<template_workspace>/exports/<template_id>_template_preview.pptx"
        ),
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Replace an existing review PPTX after an intentional re-export.",
    )
    parser.add_argument(
        "--native-charts-and-tables",
        action="store_true",
        help=(
            "Replace eligible SVG chart/table fallbacks with PowerPoint-native "
            "objects. Default review export keeps the visible SVG fallbacks."
        ),
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    require_skill_integrity()
    parser = build_parser()
    args = parser.parse_args(argv)

    try:
        workspace, template_dir = _resolve_workspace(Path(args.template_workspace))
        all_svg_files = sorted(template_dir.glob("*.svg"))
        if not all_svg_files:
            raise ValueError(f"template directory has no SVG prototypes: {template_dir}")
        definition_only_files = [
            path.name
            for path in all_svg_files
            if path.stem.startswith("layout_")
        ]
        if definition_only_files:
            raise ValueError(
                "template workspaces accept only complete Slide prototypes; "
                "replace definition-only Layout SVG(s) with authored Slide "
                "prototypes: " + ", ".join(definition_only_files)
            )
        svg_files = all_svg_files

        spec_path = _roster_spec(template_dir)
        template_id = _template_id(spec_path)
        replication_mode = _replication_mode(spec_path)
        source_themes = load_template_source_themes(template_dir)
        if source_themes is not None and replication_mode != "mirror":
            raise ValueError(
                "source_themes.json is allowed only in a mirror workspace"
            )
        locked_canvas = _canvas_viewbox(spec_path)
        if locked_canvas is None:
            raise ValueError(
                "design_spec.md frontmatter must declare canvas_viewbox"
            )
        use_full_placeholder_frames = replication_mode != "mirror"
        output_path = (
            Path(args.output).expanduser().resolve()
            if args.output
            else workspace / "exports" / f"{template_id}_template_preview.pptx"
        )
        if output_path.suffix.lower() != ".pptx":
            raise ValueError(f"output must use a .pptx extension: {output_path}")
        if output_path.exists() and not args.force:
            raise ValueError(
                f"output already exists: {output_path}; use --force to replace it"
            )
        output_path.parent.mkdir(parents=True, exist_ok=True)
        text_style, title_px, body_px = infer_master_text_style_spec(
            all_svg_files
        )

        print("PPT Master - Template Preview PPTX Exporter")
        print(f"  Workspace: {workspace}")
        print(f"  Template source: {template_dir}")
        print(f"  Slide SVG prototypes: {len(svg_files)}")
        if replication_mode == "mirror":
            print("  Review placeholder frames: preserved source Slide geometry")
        else:
            print(f"  Review Master defaults: title {title_px:g}px, body {body_px:g}px")
            print("  Review placeholder frames: full Layout bounds")
        print(f"  Output: {output_path}")

        with _review_svg_sources(
            workspace,
            svg_files,
            shorten_placeholder_markers=use_full_placeholder_frames,
        ) as review_svg_files:
            success = create_pptx_with_native_svg(
                svg_files=review_svg_files,
                output_path=output_path,
                resource_root=workspace,
                canvas_format=None,
                expected_viewbox=locked_canvas,
                verbose=True,
                transition=None,
                enable_notes=False,
                animation=None,
                image_optimize=False,
                native_objects=args.native_charts_and_tables,
                pptx_structure="structured",
                use_layout_placeholder_frames=use_full_placeholder_frames,
                master_text_style_spec=text_style,
                structure_name=template_id,
                source_theme_xml_by_master=source_themes,
            )
        if not success or not output_path.is_file():
            print("Error: template preview export did not produce a PPTX", file=sys.stderr)
            return 1

        slide_count, master_count, layout_count, placeholder_count = _verify_output(
            output_path,
            require_full_placeholder_frames=use_full_placeholder_frames,
        )
        if slide_count != len(svg_files):
            print(
                "Error: review PPTX slide count does not match the template SVG roster "
                f"({slide_count} != {len(svg_files)})",
                file=sys.stderr,
            )
            return 1

        placeholder_status = (
            f", {placeholder_count} full-frame placeholder(s)"
            if use_full_placeholder_frames
            else ""
        )
        print(
            "[OK] Template preview verified: "
            f"{slide_count} slides, {master_count} master(s), "
            f"{layout_count} layout(s){placeholder_status}"
        )
        print(output_path)
        return 0
    except (OSError, ET.ParseError, RuntimeError, ValueError) as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
